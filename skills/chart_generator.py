# -*- coding: utf-8 -*-
"""
Skill 3: chart_generator.py
PBI AutoGov -- Chart Generator

Generates chart visuals from DAX query tabular data that visually resemble
the original Power BI visuals. Supports two output formats:
  - PPTX (default): single-slide .pptx with native editable chart or PNG fallback
  - PNG (legacy): plotly-rendered static image

Input:  CSV file with DAX query results + metadata Excel from Skill 1
Output: .pptx file (one slide per visual) or .png image

Usage:
    python skills/chart_generator.py \
      --csv "output/revenue_data.csv" \
      --metadata "output/pbi_report_metadata.xlsx" \
      --visual "Pipeline by Stage" \
      --output "output/charts/"
"""

import argparse
import os
import re
import sys
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

import pandas as pd
import plotly.graph_objects as go
from plotly.subplots import make_subplots

from pptx import Presentation as _PptxFactory
from pptx.presentation import Presentation as PptxPresentation  # actual class for isinstance
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.util import Inches, Pt, Emu
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree


# =============================================================================
# CONSTANTS
# =============================================================================

# Power BI default color palette (10 accent colors)
PBI_COLORS = [
    "#118DFF", "#12239E", "#E66C37", "#6B007B", "#E044A7",
    "#744EC2", "#D9B300", "#D64550", "#197278", "#1AAB40",
]

PBI_FONT = "Segoe UI"

# PBI colors as RGBColor tuples for python-pptx native charts
PBI_RGB_COLORS = [
    RGBColor(0x11, 0x8D, 0xFF),  # #118DFF
    RGBColor(0x12, 0x23, 0x9E),  # #12239E
    RGBColor(0xE6, 0x6C, 0x37),  # #E66C37
    RGBColor(0x6B, 0x00, 0x7B),  # #6B007B
    RGBColor(0xE0, 0x44, 0xA7),  # #E044A7
    RGBColor(0x74, 0x4E, 0xC2),  # #744EC2
    RGBColor(0xD9, 0xB3, 0x00),  # #D9B300
    RGBColor(0xD6, 0x45, 0x50),  # #D64550
    RGBColor(0x19, 0x72, 0x78),  # #197278
    RGBColor(0x1A, 0xAB, 0x40),  # #1AAB40
]

# Slide dimensions: 16:9 widescreen
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Chart placement: centered with margins for title
CHART_LEFT = Inches(0.5)
CHART_TOP = Inches(0.8)
CHART_WIDTH = Inches(12.333)
CHART_HEIGHT = Inches(6.2)

# Maps PBI visual type -> XL_CHART_TYPE enum for native python-pptx rendering
NATIVE_CHART_MAP = {
    # Bar charts (horizontal)
    "barChart": XL_CHART_TYPE.BAR_CLUSTERED,
    "clusteredBarChart": XL_CHART_TYPE.BAR_CLUSTERED,
    "stackedBarChart": XL_CHART_TYPE.BAR_STACKED,
    "hundredPercentStackedBarChart": XL_CHART_TYPE.BAR_STACKED_100,
    # Column charts (vertical)
    "columnChart": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "clusteredColumnChart": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stackedColumnChart": XL_CHART_TYPE.COLUMN_STACKED,
    "hundredPercentStackedColumnChart": XL_CHART_TYPE.COLUMN_STACKED_100,
    # Line
    "lineChart": XL_CHART_TYPE.LINE_MARKERS,
    # Area
    "areaChart": XL_CHART_TYPE.AREA,
    "stackedAreaChart": XL_CHART_TYPE.AREA_STACKED,
    # Pie and donut
    "pieChart": XL_CHART_TYPE.PIE,
    "donutChart": XL_CHART_TYPE.DOUGHNUT,
    # Scatter
    "scatterChart": XL_CHART_TYPE.XY_SCATTER,
}

# Visual types that require plotly PNG fallback (no native python-pptx support)
PNG_FALLBACK_TYPES = {
    "waterfallChart", "funnelChart", "treemap", "gauge",
}


# =============================================================================
# DATA CLASS
# =============================================================================

@dataclass
class VisualSpec:
    """Metadata describing a single PBI visual for chart generation."""
    page_name: str
    visual_name: str
    visual_type: str                          # PBI camelCase identifier (e.g., "barChart")
    grouping_columns: list = field(default_factory=list)  # X-axis / Matrix Rows field names
    measure_columns: list = field(default_factory=list)   # Y-axis / Values field names
    y2_columns: list = field(default_factory=list)        # secondary-axis measures (Y2-axis)
    series_columns: list = field(default_factory=list)    # Legend / Series field names
    facet_column: str = ""                    # Small Multiples field name (if present)
    dax_pattern: str = ""                     # informational (e.g., "Pattern 3")


# =============================================================================
# PBI THEME
# =============================================================================

def get_pbi_plotly_layout():
    """Return a plotly layout dict matching PBI's default visual style.

    Applied to all chart types for consistent appearance:
    - PBI color palette as colorway
    - Segoe UI font
    - White background
    - Light gray horizontal gridlines only
    """
    return {
        "font": {"family": PBI_FONT, "size": 12, "color": "#333333"},
        "plot_bgcolor": "white",
        "paper_bgcolor": "white",
        "colorway": PBI_COLORS,
        "margin": {"l": 60, "r": 30, "t": 50, "b": 60},
        "xaxis": {
            "showgrid": False,
            "linecolor": "#E0E0E0",
            "tickfont": {"size": 10, "color": "#666666"},
        },
        "yaxis": {
            "showgrid": True,
            "gridcolor": "#E0E0E0",
            "linecolor": "#E0E0E0",
            "tickfont": {"size": 10, "color": "#666666"},
        },
    }


# =============================================================================
# DATA HELPERS
# =============================================================================

def _bare_column_name(name):
    """Extract the bare column name from a DAX-style reference.

    Handles formats like:
      - "Category[Channel]"  -> "Channel"
      - "'Category'[Channel]" -> "Channel"
      - "Channel"            -> "Channel"
    """
    import re
    m = re.search(r'\[([^\]]+)\]', name)
    return m.group(1) if m else name


def classify_columns(df, spec):
    """Split DataFrame columns into categories (grouping) and values (measures).

    Uses the VisualSpec's field lists to match against actual DataFrame column
    names (case-insensitive). Also tries matching the bare bracketed name from
    DAX-style column references (e.g. "Category[Channel]" -> "Channel").
    Falls back to dtype inference if no match.

    Returns:
        (categories: list[str], values: list[str]) — column names in df
    """
    df_cols_lower = {c.lower().strip(): c for c in df.columns}
    # Also build an index keyed by the bare bracket-extracted name
    df_cols_bare = {_bare_column_name(c).lower().strip(): c for c in df.columns}

    def _find_col(name):
        key = name.lower().strip()
        return (df_cols_lower.get(key)
                or df_cols_bare.get(key)
                or df_cols_lower.get(_bare_column_name(name).lower().strip()))

    categories = []
    for gc in spec.grouping_columns:
        actual = _find_col(gc)
        if actual and actual not in categories:
            categories.append(actual)

    values = []
    for mc in spec.measure_columns:
        actual = _find_col(mc)
        if actual and actual not in values:
            values.append(actual)

    # Fallback: infer from data types if spec columns didn't match
    if not categories and not values:
        for col in df.columns:
            if pd.api.types.is_numeric_dtype(df[col]):
                values.append(col)
            else:
                categories.append(col)

    return categories, values


def _resolve_series(df, spec):
    """Resolve spec.series_columns to actual DataFrame column names.

    Returns a list of matched column names (empty list if none found or spec has
    no series_columns). Used to pass explicit Legend/Series info to pivot helpers.
    """
    series_cols = getattr(spec, "series_columns", [])
    if not series_cols:
        return []
    df_cols_lower = {c.lower().strip(): c for c in df.columns}
    df_cols_bare = {_bare_column_name(c).lower().strip(): c for c in df.columns}

    def _find(name):
        key = name.lower().strip()
        return (df_cols_lower.get(key)
                or df_cols_bare.get(key)
                or df_cols_lower.get(_bare_column_name(name).lower().strip()))

    resolved = []
    for sc in series_cols:
        actual = _find(sc)
        if actual and actual not in resolved:
            resolved.append(actual)
    return resolved


_MONTH_ORDER = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
                "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
_MONTH_RANK = {m: i for i, m in enumerate(_MONTH_ORDER)}


def _sort_categories(labels):
    """Sort category labels chronologically if they look like month abbreviations,
    otherwise return as-is."""
    if all(str(l).strip() in _MONTH_RANK for l in labels):
        return sorted(labels, key=lambda m: _MONTH_RANK[str(m).strip()])
    return labels


def _prepare_series_data(df, categories, values, series=None):
    """Prepare category labels and series data for bar/column/stacked charts.

    Pivot logic priority:
    1. If `series` is provided (explicit Legend well field), pivot on that column.
    2. If 2+ grouping columns and 1 measure, pivot on the second grouping (legacy).
    3. Otherwise, each measure becomes its own series.

    Args:
        series: list of df column names to use as the legend/series pivot axis.
                When provided (well-aware mode), overrides the legacy 2-grouping heuristic.

    Returns:
        (cat_labels: list[str], series_data: OrderedDict[str, list[float]], needs_legend: bool)
    """
    from collections import OrderedDict

    # Well-aware: explicit Legend column provided
    if series and len(series) == 1 and len(values) == 1 and categories:
        pivot_df = df.pivot_table(
            index=categories[0], columns=series[0],
            values=values[0], aggfunc="sum"
        ).fillna(0)
        sorted_index = _sort_categories(list(pivot_df.index))
        pivot_df = pivot_df.reindex(sorted_index)
        cat_labels = [str(c) for c in pivot_df.index]
        series_data = OrderedDict((str(col), pivot_df[col].tolist()) for col in pivot_df.columns)
        return cat_labels, series_data, True

    # Legacy: treat second grouping column as legend when no explicit series
    if len(categories) >= 2 and len(values) == 1:
        pivot_df = df.pivot_table(
            index=categories[0], columns=categories[1],
            values=values[0], aggfunc="sum"
        ).fillna(0)
        sorted_index = _sort_categories(list(pivot_df.index))
        pivot_df = pivot_df.reindex(sorted_index)
        cat_labels = [str(c) for c in pivot_df.index]
        series_data = OrderedDict((str(col), pivot_df[col].tolist()) for col in pivot_df.columns)
        return cat_labels, series_data, True

    sorted_vals = _sort_categories(df[categories[0]].astype(str).tolist())
    if sorted_vals != df[categories[0]].astype(str).tolist():
        df = df.set_index(categories[0]).reindex(sorted_vals).reset_index()
    cat_labels = df[categories[0]].astype(str).tolist()
    series_data = OrderedDict((v, df[v].tolist()) for v in values)
    needs_legend = len(values) > 1
    return cat_labels, series_data, needs_legend


def parse_visual_from_metadata(metadata_excel, visual_name):
    """Read metadata Excel (from Skill 1) and build a VisualSpec for a named visual.

    Imports classify_field from dax_query_builder to determine field roles.
    Matches visual_name case-insensitively (partial match supported).

    Returns:
        VisualSpec or None if visual not found
    """
    # Import from sibling module (same directory)
    skills_dir = os.path.dirname(os.path.abspath(__file__))
    if skills_dir not in sys.path:
        sys.path.insert(0, skills_dir)
    from dax_query_builder import read_extractor_output, classify_field

    visuals, _, _, _ = read_extractor_output(metadata_excel)

    # Find the visual by name (case-insensitive, supports "Page / Visual" format)
    target = visual_name.lower().strip()
    matched_key = None

    # First try exact match on "page / visual" combined key
    for key, data in visuals.items():
        page = key[0]
        full_name = f"{page} / {data['visual_name']}".lower()
        if full_name == target:
            matched_key = key
            break

    # Then try exact match on visual name alone
    if not matched_key:
        for key, data in visuals.items():
            if data["visual_name"].lower().strip() == target:
                matched_key = key
                break

    # Partial match fallback (also checks "page / visual" combined)
    if not matched_key:
        for key, data in visuals.items():
            page = key[0]
            full_name = f"{page} / {data['visual_name']}".lower()
            if target in data["visual_name"].lower() or target in full_name:
                matched_key = key
                break

    if not matched_key:
        print(f"ERROR: Visual '{visual_name}' not found in metadata.")
        print(f"Available visuals:")
        for key, data in visuals.items():
            print(f"  - {data['visual_name']} ({data['visual_type']}) on page '{key[0]}'")
        return None

    data = visuals[matched_key]
    page_name = matched_key[0]

    # Check whether the Excel has well assignments (new metadata) or only usage labels (legacy)
    has_well = any(f.get("well", "") for f in data["fields"])

    grouping_cols = []    # X-axis / Matrix Rows
    series_cols = []      # Legend / Series
    facet_col = ""        # Small Multiples
    measure_cols = []     # Y-axis / Values
    y2_cols = []          # Y2-axis
    seen_field = set()    # deduplicate by (table, col) key

    for f in data["fields"]:
        well = f.get("well", "")
        ui_name = f["ui_name"]
        field_key = (f["table_sm"], f["col_sm"])

        if has_well and well:
            w = well.lower()
            if w == "small multiples":
                if not facet_col:
                    facet_col = ui_name
            elif w == "legend":
                if ui_name not in series_cols:
                    series_cols.append(ui_name)
            elif w in ("x-axis", "matrix rows"):
                if field_key not in seen_field:
                    seen_field.add(field_key)
                    grouping_cols.append(ui_name)
            elif w == "y2-axis":
                if ui_name not in measure_cols:
                    measure_cols.append(ui_name)
                if ui_name not in y2_cols:
                    y2_cols.append(ui_name)
            elif w in ("y-axis", "values", "indicator", "trend", "target",
                       "target value", "size", "tooltip", "analyze"):
                if ui_name not in measure_cols:
                    measure_cols.append(ui_name)
            # Ignore wells like "Tooltip", "Details", "Play Axis" for chart layout
        else:
            # Legacy fallback: classify by usage label
            role = classify_field(f["usage"])
            if role == "grouping":
                if field_key not in seen_field:
                    seen_field.add(field_key)
                    grouping_cols.append(ui_name)
            elif role == "measure":
                if ui_name not in measure_cols:
                    measure_cols.append(ui_name)
                if "y2" in f["usage"].lower() and ui_name not in y2_cols:
                    y2_cols.append(ui_name)

    return VisualSpec(
        page_name=page_name,
        visual_name=data["visual_name"],
        visual_type=data["visual_type"],
        grouping_columns=grouping_cols,
        measure_columns=measure_cols,
        y2_columns=y2_cols,
        series_columns=series_cols,
        facet_column=facet_col,
    )


# =============================================================================
# CHART RENDERERS — each returns a plotly Figure
# =============================================================================

# ---- Bar chart (horizontal) ----

def _render_bar(df, spec):
    """Horizontal bar chart. Categories on Y-axis, values on X-axis.

    Handles single measure (one series) or multiple measures (grouped bars).
    With 2+ grouping columns, pivots second column into legend series.
    """
    categories, values = classify_columns(df, spec)
    if not categories or not values:
        return _render_table(df, spec)

    # Sort by first measure descending for clean ranking (like PBI default)
    if len(values) == 1 and len(categories) == 1:
        df = df.sort_values(values[0], ascending=False)  # highest at top with reversed y-axis

    series_cols = _resolve_series(df, spec)
    cat_labels, series, needs_legend = _prepare_series_data(df, categories, values, series_cols)

    fig = go.Figure()
    for i, (name, data) in enumerate(series.items()):
        fig.add_trace(go.Bar(
            y=cat_labels, x=data, name=name, orientation="h",
            marker_color=PBI_COLORS[i % len(PBI_COLORS)],
            text=[f"{v:,.0f}" for v in data],
            textposition="outside",
            textfont={"size": 9, "family": PBI_FONT},
        ))

    fig.update_layout(
        **get_pbi_plotly_layout(),
        title=spec.visual_name,
        barmode="group",
        showlegend=needs_legend,
        xaxis_title=values[0] if len(values) == 1 else None,
        yaxis_title=categories[0] if categories else None,
    )
    fig.update_yaxes(autorange="reversed")  # top-to-bottom ordering like PBI
    return fig


# ---- Column chart (vertical) ----

def _render_column(df, spec):
    """Vertical column chart. Categories on X-axis, values on Y-axis."""
    categories, values = classify_columns(df, spec)
    if not categories or not values:
        return _render_table(df, spec)

    # Sort by first measure descending
    if len(values) == 1 and len(categories) == 1:
        df = df.sort_values(values[0], ascending=False)

    series_cols = _resolve_series(df, spec)
    cat_labels, series, needs_legend = _prepare_series_data(df, categories, values, series_cols)

    fig = go.Figure()
    for i, (name, data) in enumerate(series.items()):
        fig.add_trace(go.Bar(
            x=cat_labels, y=data, name=name,
            marker_color=PBI_COLORS[i % len(PBI_COLORS)],
            text=[f"{v:,.0f}" for v in data],
            textposition="outside",
            textfont={"size": 9, "family": PBI_FONT},
        ))

    fig.update_layout(
        **get_pbi_plotly_layout(),
        title=spec.visual_name,
        barmode="group",
        showlegend=needs_legend,
        xaxis_title=categories[0] if categories else None,
        yaxis_title=values[0] if len(values) == 1 else None,
    )
    return fig


# ---- Stacked bar (horizontal) ----

def _render_stacked_bar(df, spec):
    """Horizontal stacked bar chart. Uses barmode='stack' or barnorm='percent'."""
    categories, values = classify_columns(df, spec)
    if not categories or not values:
        return _render_table(df, spec)

    series_cols = _resolve_series(df, spec)
    cat_labels, series, _ = _prepare_series_data(df, categories, values, series_cols)

    fig = go.Figure()
    for i, (name, data) in enumerate(series.items()):
        fig.add_trace(go.Bar(
            y=cat_labels, x=data, name=name, orientation="h",
            marker_color=PBI_COLORS[i % len(PBI_COLORS)],
            text=[f"{v:,.0f}" for v in data],
            textposition="inside",
            textfont={"size": 9, "family": PBI_FONT, "color": "white"},
        ))

    layout_kwargs = {
        **get_pbi_plotly_layout(),
        "title": spec.visual_name,
        "barmode": "stack",
        "showlegend": True,
        "xaxis_title": values[0] if len(values) == 1 else None,
        "yaxis_title": categories[0] if categories else None,
    }
    if "hundredPercent" in spec.visual_type:
        layout_kwargs["barnorm"] = "percent"
    fig.update_layout(**layout_kwargs)
    fig.update_yaxes(autorange="reversed")  # top-to-bottom ordering like PBI
    return fig


# ---- Stacked column (vertical) ----

def _render_stacked_column(df, spec):
    """Vertical stacked column chart. Uses barmode='stack' or barnorm='percent'."""
    categories, values = classify_columns(df, spec)
    if not categories or not values:
        return _render_table(df, spec)

    series_cols = _resolve_series(df, spec)
    cat_labels, series, _ = _prepare_series_data(df, categories, values, series_cols)

    fig = go.Figure()
    for i, (name, data) in enumerate(series.items()):
        fig.add_trace(go.Bar(
            x=cat_labels, y=data, name=name,
            marker_color=PBI_COLORS[i % len(PBI_COLORS)],
            text=[f"{v:,.0f}" for v in data],
            textposition="inside",
            textfont={"size": 9, "family": PBI_FONT, "color": "white"},
        ))

    layout_kwargs = {
        **get_pbi_plotly_layout(),
        "title": spec.visual_name,
        "barmode": "stack",
        "showlegend": True,
        "xaxis_title": categories[0] if categories else None,
        "yaxis_title": values[0] if len(values) == 1 else None,
    }
    if "hundredPercent" in spec.visual_type:
        layout_kwargs["barnorm"] = "percent"
    fig.update_layout(**layout_kwargs)
    return fig


# ---- Line chart ----

def _render_line(df, spec):
    """Line chart. X-axis is first grouping column, one line per measure.

    Data is sorted by the category column for proper line ordering.
    With 2+ grouping columns, pivots second into legend series.
    """
    categories, values = classify_columns(df, spec)
    if not categories or not values:
        return _render_table(df, spec)

    df_sorted = df.sort_values(categories[0])

    if len(categories) >= 2 and len(values) == 1:
        pivot_df = df_sorted.pivot_table(
            index=categories[0], columns=categories[1],
            values=values[0], aggfunc="sum"
        ).fillna(0)
        fig = go.Figure()
        for i, col in enumerate(pivot_df.columns):
            col_data = pivot_df[col].tolist()
            fig.add_trace(go.Scatter(
                x=pivot_df.index.astype(str), y=col_data,
                name=str(col), mode="lines+markers+text",
                line={"color": PBI_COLORS[i % len(PBI_COLORS)]},
                text=[f"{v:,.0f}" for v in col_data],
                textposition="top center",
                textfont={"size": 9, "family": PBI_FONT},
            ))
        needs_legend = True
    else:
        fig = go.Figure()
        cat_labels = df_sorted[categories[0]].astype(str).tolist()
        for i, v in enumerate(values):
            v_data = df_sorted[v].tolist()
            fig.add_trace(go.Scatter(
                x=cat_labels, y=v_data,
                name=v, mode="lines+markers+text",
                line={"color": PBI_COLORS[i % len(PBI_COLORS)]},
                text=[f"{val:,.0f}" for val in v_data],
                textposition="top center",
                textfont={"size": 9, "family": PBI_FONT},
            ))
        needs_legend = len(values) > 1

    fig.update_layout(
        **get_pbi_plotly_layout(),
        title=spec.visual_name,
        showlegend=needs_legend,
        xaxis_title=categories[0] if categories else None,
        yaxis_title=values[0] if len(values) == 1 else None,
    )
    return fig


# ---- Area chart ----

def _render_area(df, spec):
    """Area chart. Like line chart but with fill to zero.

    Stacked area uses stackgroup for overlapping fills.
    """
    categories, values = classify_columns(df, spec)
    if not categories or not values:
        return _render_table(df, spec)

    is_stacked = "stacked" in spec.visual_type.lower()
    df_sorted = df.sort_values(categories[0])

    fig = go.Figure()
    cat_labels = df_sorted[categories[0]].astype(str).tolist()
    for i, v in enumerate(values):
        trace_kwargs = {
            "x": cat_labels,
            "y": df_sorted[v].tolist(),
            "name": v,
            "mode": "lines",
            "line": {"color": PBI_COLORS[i % len(PBI_COLORS)]},
        }
        if is_stacked:
            trace_kwargs["stackgroup"] = "one"
        else:
            trace_kwargs["fill"] = "tozeroy"
        fig.add_trace(go.Scatter(**trace_kwargs))

    fig.update_layout(
        **get_pbi_plotly_layout(),
        title=spec.visual_name,
        showlegend=len(values) > 1,
        xaxis_title=categories[0] if categories else None,
        yaxis_title=values[0] if len(values) == 1 else None,
    )
    return fig


# ---- Pie chart ----

def _render_pie(df, spec):
    """Pie chart. Labels from first grouping column, values from first measure.

    Only uses first measure (pie charts show one value series). Warns if multiple.
    Measures-only case: each measure name becomes a slice label, each value a slice size.
    """
    categories, values = classify_columns(df, spec)

    # Measures-only: no grouping column, each measure is a slice
    if not categories and len(values) >= 2:
        row = df.iloc[0]
        slice_labels = values
        slice_values = [float(row[v]) if pd.notna(row[v]) else 0 for v in values]
        fig = go.Figure(go.Pie(
            labels=slice_labels,
            values=slice_values,
            marker={"colors": PBI_COLORS[:len(values)]},
            textinfo="percent+label",
            textfont={"size": 11, "family": PBI_FONT},
            hole=0,
        ))
        fig.update_layout(
            title=spec.visual_name,
            font={"family": PBI_FONT, "size": 12, "color": "#333333"},
            paper_bgcolor="white",
            showlegend=True,
            legend={"font": {"size": 10}},
        )
        return fig

    if not categories or not values:
        return _render_table(df, spec)

    if len(values) > 1:
        print(f"  WARNING: Pie chart '{spec.visual_name}' has {len(values)} measures "
              f"-- using only '{values[0]}'")

    fig = go.Figure(go.Pie(
        labels=df[categories[0]].astype(str).tolist(),
        values=df[values[0]].tolist(),
        marker={"colors": PBI_COLORS[:len(df)]},
        textinfo="percent+label+value",
        textfont={"size": 11, "family": PBI_FONT},
        hole=0,
    ))
    fig.update_layout(
        title=spec.visual_name,
        font={"family": PBI_FONT, "size": 12, "color": "#333333"},
        paper_bgcolor="white",
        showlegend=True,
        legend={"font": {"size": 10}},
    )
    return fig


# ---- Donut chart ----

def _render_donut(df, spec):
    """Donut chart. Same as pie but with a hole in the center (hole=0.4).

    Measures-only case: each measure name becomes a slice label, each value a slice size.
    """
    categories, values = classify_columns(df, spec)

    # Measures-only: no grouping column, each measure is a slice
    if not categories and len(values) >= 2:
        row = df.iloc[0]
        slice_labels = values
        slice_values = [float(row[v]) if pd.notna(row[v]) else 0 for v in values]
        fig = go.Figure(go.Pie(
            labels=slice_labels,
            values=slice_values,
            marker={"colors": PBI_COLORS[:len(values)]},
            textinfo="percent+label+value",
            textfont={"size": 11, "family": PBI_FONT},
            hole=0.4,
        ))
        fig.update_layout(
            title=spec.visual_name,
            font={"family": PBI_FONT, "size": 12, "color": "#333333"},
            paper_bgcolor="white",
            showlegend=True,
            legend={"font": {"size": 10}},
        )
        return fig

    if not categories or not values:
        return _render_table(df, spec)

    if len(values) > 1:
        print(f"  WARNING: Donut chart '{spec.visual_name}' has {len(values)} measures "
              f"-- using only '{values[0]}'")

    fig = go.Figure(go.Pie(
        labels=df[categories[0]].astype(str).tolist(),
        values=df[values[0]].tolist(),
        marker={"colors": PBI_COLORS[:len(df)]},
        textinfo="percent+label+value",
        textfont={"size": 11, "family": PBI_FONT},
        hole=0.4,
    ))
    fig.update_layout(
        title=spec.visual_name,
        font={"family": PBI_FONT, "size": 12, "color": "#333333"},
        paper_bgcolor="white",
        showlegend=True,
        legend={"font": {"size": 10}},
    )
    return fig


# ---- Scatter chart ----

def _render_scatter(df, spec):
    """Scatter plot. X = first measure, Y = second measure.

    If a grouping column exists, it becomes the point labels/color grouper.
    Third measure (if present) maps to marker size (bubble chart).
    """
    categories, values = classify_columns(df, spec)
    if len(values) < 2:
        return _render_table(df, spec)

    # Guard for bubble sizing: ensure positive max for sizeref calculation
    has_bubble = len(values) >= 3
    if has_bubble:
        max_bubble = df[values[2]].max()
        # Fallback to fixed size if all bubble values are zero or negative
        if pd.isna(max_bubble) or max_bubble <= 0:
            has_bubble = False

    fig = go.Figure()

    if categories:
        groups = df.groupby(categories[0])
        for i, (name, group) in enumerate(groups):
            trace_kwargs = {
                "x": group[values[0]].tolist(),
                "y": group[values[1]].tolist(),
                "name": str(name),
                "mode": "markers+text",
                "marker": {"color": PBI_COLORS[i % len(PBI_COLORS)], "size": 10},
                "text": group[categories[0]].astype(str).tolist(),
                "textposition": "top center",
                "textfont": {"size": 9, "family": PBI_FONT},
            }
            if has_bubble:
                trace_kwargs["marker"]["size"] = group[values[2]].tolist()
                trace_kwargs["marker"]["sizemode"] = "area"
                trace_kwargs["marker"]["sizeref"] = 2.0 * max_bubble / (40.0 ** 2)
            fig.add_trace(go.Scatter(**trace_kwargs))
    else:
        y_data = df[values[1]].tolist()
        trace_kwargs = {
            "x": df[values[0]].tolist(),
            "y": y_data,
            "mode": "markers+text",
            "marker": {"color": PBI_COLORS[0], "size": 10},
            "text": [f"{v:,.0f}" for v in y_data],
            "textposition": "top center",
            "textfont": {"size": 9, "family": PBI_FONT},
        }
        if has_bubble:
            trace_kwargs["marker"]["size"] = df[values[2]].tolist()
            trace_kwargs["marker"]["sizemode"] = "area"
            trace_kwargs["marker"]["sizeref"] = 2.0 * max_bubble / (40.0 ** 2)
        fig.add_trace(go.Scatter(**trace_kwargs))

    fig.update_layout(
        **get_pbi_plotly_layout(),
        title=spec.visual_name,
        xaxis_title=values[0],
        yaxis_title=values[1],
        showlegend=bool(categories),
    )
    return fig


# ---- Waterfall chart ----

def _render_waterfall(df, spec):
    """Waterfall chart using plotly's native go.Waterfall trace.

    Categories from first grouping column, values from first measure.
    All values treated as relative (incremental). PBI blue for positive,
    PBI red for negative.
    """
    categories, values = classify_columns(df, spec)
    if not categories or not values:
        return _render_table(df, spec)

    cat_labels = df[categories[0]].astype(str).tolist()
    val_data = df[values[0]].tolist()

    # All bars are relative (incremental values)
    measure_types = ["relative"] * len(val_data)

    fig = go.Figure(go.Waterfall(
        x=cat_labels,
        y=val_data,
        measure=measure_types,
        connector={"line": {"color": "#E0E0E0"}},
        increasing={"marker": {"color": PBI_COLORS[0]}},   # blue for positive
        decreasing={"marker": {"color": PBI_COLORS[7]}},    # red for negative
        totals={"marker": {"color": PBI_COLORS[1]}},        # dark blue for totals
        text=[f"{v:,.0f}" for v in val_data],
        textposition="outside",
        textfont={"size": 9, "family": PBI_FONT},
    ))
    fig.update_layout(
        **get_pbi_plotly_layout(),
        title=spec.visual_name,
        showlegend=False,
        xaxis_title=categories[0] if categories else None,
        yaxis_title=values[0] if len(values) == 1 else None,
    )
    return fig


# ---- Combo chart (dual axis: bars + line) ----

def _render_combo(df, spec):
    """Combo chart renderer.

    Two modes:
    1. Small multiples (single measure + 2 grouping columns): renders one
       column subplot per unique value of the first grouping (e.g. Channel),
       with the second grouping (e.g. Months) as the X-axis. This matches
       the PBI "small multiples" layout.
    2. Dual-axis combo (2+ measures): bars on primary Y, line on secondary Y.
    """
    categories, values = classify_columns(df, spec)
    if not categories or not values:
        return _render_table(df, spec)

    # --- Mode 1: small multiples ---
    # Use spec.facet_column (from Small Multiples well) if available,
    # otherwise fall back to the legacy 2-grouping-columns heuristic.
    explicit_facet = getattr(spec, "facet_column", "")
    if explicit_facet or (len(values) == 1 and len(categories) >= 2):
        # Resolve facet column from df
        df_cols_lower_combo = {c.lower().strip(): c for c in df.columns}
        df_cols_bare_combo = {_bare_column_name(c).lower().strip(): c for c in df.columns}
        def _find_col_combo(name):
            key = name.lower().strip()
            return (df_cols_lower_combo.get(key)
                    or df_cols_bare_combo.get(key)
                    or df_cols_lower_combo.get(_bare_column_name(name).lower().strip()))

        if explicit_facet:
            facet_col_actual = _find_col_combo(explicit_facet)
        else:
            facet_col_actual = None

        # Determine x_col: the first grouping column that isn't the facet
        if facet_col_actual and categories:
            x_col_actual = categories[0]
        elif len(categories) >= 2:
            facet_col_actual = categories[0]
            x_col_actual = categories[1]
        else:
            return _render_table(df, spec)

        facet_col = facet_col_actual
        x_col = x_col_actual
        measure = values[0]

        facet_values = df[facet_col].dropna().unique().tolist()
        # Preserve a natural order if possible (sort alphabetically as fallback)
        facet_values = sorted(facet_values, key=str)
        n_facets = len(facet_values)

        # Chronological sort for x-axis
        all_x = _sort_categories(df[x_col].dropna().unique().tolist())

        fig = make_subplots(
            rows=1, cols=n_facets,
            subplot_titles=[str(f) for f in facet_values],
            shared_yaxes=True,
        )

        for col_idx, facet_val in enumerate(facet_values, start=1):
            sub_df = df[df[facet_col] == facet_val].copy()
            # Reindex to full x_order, fill missing with 0
            sub_df = sub_df.set_index(x_col).reindex(all_x).fillna(0).reset_index()
            y_vals = sub_df[measure].tolist()
            color = PBI_COLORS[col_idx - 1]

            fig.add_trace(
                go.Bar(
                    x=all_x,
                    y=y_vals,
                    name=str(facet_val),
                    marker_color=color,
                    showlegend=False,
                    text=[f"${v/1000:.0f}K" if v >= 1000 else f"${v:.0f}" for v in y_vals],
                    textposition="outside",
                    textfont={"size": 8, "family": PBI_FONT},
                ),
                row=1, col=col_idx,
            )
            # Subtitle annotation (facet label at the bottom)
            fig.update_xaxes(
                tickangle=45,
                tickfont={"size": 9, "family": PBI_FONT},
                showgrid=False,
                linecolor="#E0E0E0",
                row=1, col=col_idx,
            )
            fig.update_yaxes(
                showgrid=True,
                gridcolor="#E0E0E0",
                tickformat="$,.0f",
                tickfont={"size": 9, "family": PBI_FONT},
                row=1, col=col_idx,
            )

        fig.update_layout(
            title=spec.visual_name,
            font={"family": PBI_FONT, "size": 11, "color": "#333333"},
            plot_bgcolor="white",
            paper_bgcolor="white",
            bargap=0.15,
            margin={"l": 60, "r": 20, "t": 60, "b": 60},
            height=420,
        )
        return fig

    # --- Mode 2: dual-axis combo (2+ measures) ---
    sorted_index = _sort_categories(df[categories[0]].astype(str).tolist())
    if sorted_index != df[categories[0]].astype(str).tolist():
        df = df.set_index(categories[0]).reindex(sorted_index).reset_index()
    cat_labels = df[categories[0]].astype(str).tolist()

    # Resolve actual column names for y2 (case-insensitive match against df)
    df_cols_lower = {c.lower().strip(): c for c in df.columns}
    y2_actual = set()
    for y2 in spec.y2_columns:
        actual = df_cols_lower.get(y2.lower().strip())
        if actual:
            y2_actual.add(actual)

    # Split measures into bar vs line using y2 metadata when available
    if y2_actual:
        bar_measures = [v for v in values if v not in y2_actual]
        line_measures = [v for v in values if v in y2_actual]
    elif len(values) > 1:
        # Fallback: all except last → bars, last → line
        bar_measures = values[:-1]
        line_measures = [values[-1]]
    else:
        bar_measures = values
        line_measures = []

    fig = make_subplots(specs=[[{"secondary_y": bool(line_measures)}]])

    for i, m in enumerate(bar_measures):
        bar_data = df[m].tolist()
        fig.add_trace(go.Bar(
            x=cat_labels, y=bar_data, name=m,
            marker_color=PBI_COLORS[i % len(PBI_COLORS)],
            text=[f"{v:,.0f}" for v in bar_data],
            textposition="outside",
            textfont={"size": 9, "family": PBI_FONT},
        ), secondary_y=False)

    for i, m in enumerate(line_measures):
        line_data = df[m].tolist()
        fig.add_trace(go.Scatter(
            x=cat_labels, y=line_data, name=m,
            mode="lines+markers+text",
            line={"color": PBI_COLORS[(len(bar_measures) + i) % len(PBI_COLORS)], "width": 2},
            text=[f"{v:,.0f}" for v in line_data],
            textposition="top center",
            textfont={"size": 9, "family": PBI_FONT},
        ), secondary_y=True)

    layout = get_pbi_plotly_layout()
    fig.update_layout(
        title=spec.visual_name,
        font=layout["font"],
        plot_bgcolor=layout["plot_bgcolor"],
        paper_bgcolor=layout["paper_bgcolor"],
        barmode="group",
        showlegend=True,
        margin=layout["margin"],
    )
    fig.update_xaxes(showgrid=False, linecolor="#E0E0E0",
                     title_text=categories[0] if categories else None)
    fig.update_yaxes(showgrid=True, gridcolor="#E0E0E0", linecolor="#E0E0E0")
    if bar_measures:
        fig.update_yaxes(title_text=bar_measures[0] if len(bar_measures) == 1 else None,
                         secondary_y=False)
    if line_measures:
        fig.update_yaxes(title_text=line_measures[0] if len(line_measures) == 1 else None,
                         secondary_y=True)
    return fig


# ---- Funnel chart ----

def _render_funnel(df, spec):
    """Funnel chart. Categories (stages) on Y-axis, values (counts) on X-axis."""
    categories, values = classify_columns(df, spec)
    if not categories or not values:
        return _render_table(df, spec)

    fig = go.Figure(go.Funnel(
        y=df[categories[0]].astype(str).tolist(),
        x=df[values[0]].tolist(),
        marker={"color": PBI_COLORS[:len(df)]},
        textinfo="value+percent initial",
        textfont={"family": PBI_FONT, "size": 11},
    ))
    fig.update_layout(
        title=spec.visual_name,
        font={"family": PBI_FONT, "size": 12, "color": "#333333"},
        paper_bgcolor="white",
        plot_bgcolor="white",
        showlegend=False,
        yaxis_title=categories[0] if categories else None,
        xaxis_title=values[0] if len(values) == 1 else None,
    )
    return fig


# ---- Treemap ----

def _render_treemap(df, spec):
    """Treemap chart. Labels from grouping column, tile sizes from measure.

    With 2+ grouping columns, creates nested hierarchy (parent -> child).
    """
    categories, values = classify_columns(df, spec)
    if not categories or not values:
        return _render_table(df, spec)

    if len(categories) >= 2:
        # Nested treemap: first grouping = parent, second = child
        labels = df[categories[1]].astype(str).tolist()
        parents = df[categories[0]].astype(str).tolist()
    else:
        labels = df[categories[0]].astype(str).tolist()
        parents = [""] * len(labels)

    fig = go.Figure(go.Treemap(
        labels=labels,
        parents=parents,
        values=df[values[0]].tolist(),
        marker={"colors": PBI_COLORS[:len(df)]},
        textinfo="label+value",
        textfont={"family": PBI_FONT, "size": 12},
    ))
    fig.update_layout(
        title=spec.visual_name,
        font={"family": PBI_FONT, "size": 12, "color": "#333333"},
        paper_bgcolor="white",
        margin={"l": 10, "r": 10, "t": 50, "b": 10},
    )
    return fig


# ---- Gauge ----

def _render_gauge(df, spec):
    """Gauge indicator. Value from first measure, optional target from second.

    Range auto-calculated as 0 to value * 1.2.
    """
    _, values = classify_columns(df, spec)
    if not values or df.empty:
        return _render_table(df, spec)

    value = float(df[values[0]].iloc[0])
    max_range = value * 1.2 if value > 0 else 100

    gauge_kwargs = {
        "axis": {"range": [0, max_range], "tickfont": {"size": 10}},
        "bar": {"color": PBI_COLORS[0]},
        "bgcolor": "white",
        "borderwidth": 0,
        "steps": [
            {"range": [0, max_range * 0.5], "color": "#F0F0F0"},
            {"range": [max_range * 0.5, max_range], "color": "#E0E0E0"},
        ],
    }

    if len(values) >= 2:
        target = float(df[values[1]].iloc[0])
        gauge_kwargs["threshold"] = {
            "line": {"color": PBI_COLORS[7], "width": 3},
            "thickness": 0.8,
            "value": target,
        }

    fig = go.Figure(go.Indicator(
        mode="gauge+number",
        value=value,
        title={"text": spec.visual_name, "font": {"size": 16, "family": PBI_FONT}},
        number={"font": {"size": 36, "family": PBI_FONT, "color": "#333333"}},
        gauge=gauge_kwargs,
    ))
    fig.update_layout(
        paper_bgcolor="white",
        font={"family": PBI_FONT},
        margin={"l": 30, "r": 30, "t": 60, "b": 30},
    )
    return fig


# ---- Card (single or multi-value) ----

def _render_card(df, spec):
    """Card visual. Displays measure values as large centered numbers.

    Single measure -> one big number. Multiple measures -> side-by-side indicators.
    Uses go.Indicator(mode="number") for clean big-number display.
    """
    _, values = classify_columns(df, spec)
    if not values or df.empty:
        return _render_table(df, spec)

    row = df.iloc[0]

    if len(values) == 1:
        val = row[values[0]]
        fig = go.Figure(go.Indicator(
            mode="number",
            value=float(val) if pd.notna(val) else 0,
            title={"text": values[0], "font": {"size": 16, "family": PBI_FONT, "color": "#666666"}},
            number={"font": {"size": 60, "family": PBI_FONT, "color": PBI_COLORS[0]},
                    "valueformat": ",.0f"},
        ))
        fig.update_layout(
            paper_bgcolor="white",
            margin={"l": 30, "r": 30, "t": 60, "b": 30},
        )
    else:
        fig = make_subplots(
            rows=1, cols=len(values),
            specs=[[{"type": "indicator"}] * len(values)],
        )
        for i, v in enumerate(values):
            val = row[v]
            fig.add_trace(go.Indicator(
                mode="number",
                value=float(val) if pd.notna(val) else 0,
                title={"text": v, "font": {"size": 13, "family": PBI_FONT, "color": "#666666"}},
                number={"font": {"size": 36, "family": PBI_FONT, "color": PBI_COLORS[i % len(PBI_COLORS)]},
                        "valueformat": ",.0f"},
            ), row=1, col=i + 1)
        fig.update_layout(
            paper_bgcolor="white",
            margin={"l": 20, "r": 20, "t": 60, "b": 20},
        )

    return fig


# ---- KPI ----

def _render_kpi(df, spec):
    """KPI visual. Shows value with delta (change from previous/target).

    First measure = value, second measure = reference (for delta calculation).
    """
    _, values = classify_columns(df, spec)
    if not values or df.empty:
        return _render_table(df, spec)

    row = df.iloc[0]
    value = float(row[values[0]]) if pd.notna(row[values[0]]) else 0

    indicator_kwargs = {
        "mode": "number+delta",
        "value": value,
        "title": {"text": spec.visual_name, "font": {"size": 16, "family": PBI_FONT}},
        "number": {"font": {"size": 48, "family": PBI_FONT, "color": "#333333"},
                   "valueformat": ",.0f"},
    }

    if len(values) >= 2:
        reference = float(row[values[1]]) if pd.notna(row[values[1]]) else 0
        indicator_kwargs["delta"] = {
            "reference": reference,
            "relative": True,
            "valueformat": ".1%",
            "increasing": {"color": "#1AAB40"},
            "decreasing": {"color": "#D64550"},
        }

    fig = go.Figure(go.Indicator(**indicator_kwargs))
    fig.update_layout(
        paper_bgcolor="white",
        font={"family": PBI_FONT},
        margin={"l": 30, "r": 30, "t": 60, "b": 30},
    )
    return fig


# ---- Ribbon chart ----

def _render_ribbon(df, spec):
    """Ribbon chart rendered as a stacked area chart (closest plotly equivalent)."""
    categories, values = classify_columns(df, spec)
    if not categories or not values:
        return _render_table(df, spec)

    df_sorted = df.sort_values(categories[0])
    fig = go.Figure()

    if len(categories) >= 2 and len(values) == 1:
        pivot_df = df_sorted.pivot_table(
            index=categories[0], columns=categories[1],
            values=values[0], aggfunc="sum"
        ).fillna(0)
        for i, col in enumerate(pivot_df.columns):
            fig.add_trace(go.Scatter(
                x=pivot_df.index.astype(str), y=pivot_df[col],
                name=str(col), mode="lines", stackgroup="one",
                line={"color": PBI_COLORS[i % len(PBI_COLORS)]},
            ))
    else:
        cat_labels = df_sorted[categories[0]].astype(str).tolist()
        for i, v in enumerate(values):
            fig.add_trace(go.Scatter(
                x=cat_labels, y=df_sorted[v].tolist(),
                name=v, mode="lines", stackgroup="one",
                line={"color": PBI_COLORS[i % len(PBI_COLORS)]},
            ))

    fig.update_layout(
        **get_pbi_plotly_layout(),
        title=spec.visual_name,
        showlegend=True,
    )
    return fig


# ---- Table (also used as fallback for unknown types) ----

def _render_table(df, spec):
    """Table visual using go.Table. Also the fallback for unknown visual types.

    Renders up to 50 rows with PBI-styled header and alternating row colors.
    """
    max_rows = 50
    display_df = df.head(max_rows)
    num_rows = len(display_df)

    header_values = list(display_df.columns)
    cell_values = [display_df[col].astype(str).tolist() for col in display_df.columns]

    # Build row color list matching exact row count
    row_colors = ["white" if i % 2 == 0 else "#F5F5F5" for i in range(num_rows)]

    fig = go.Figure(go.Table(
        header={
            "values": [f"<b>{h}</b>" for h in header_values],
            "fill_color": PBI_COLORS[1],        # dark blue header
            "font": {"color": "white", "size": 11, "family": PBI_FONT},
            "align": "left",
            "height": 30,
        },
        cells={
            "values": cell_values,
            "fill_color": [row_colors],
            "font": {"size": 10, "family": PBI_FONT, "color": "#333333"},
            "align": "left",
            "height": 25,
        },
    ))

    title = spec.visual_name
    if len(df) > max_rows:
        title += f" (showing {max_rows} of {len(df)} rows)"

    fig.update_layout(
        title=title,
        font={"family": PBI_FONT},
        paper_bgcolor="white",
        margin={"l": 10, "r": 10, "t": 50, "b": 10},
    )
    return fig


# =============================================================================
# CHART TYPE ROUTER
# =============================================================================

# Maps PBI visual type identifiers to renderer functions
CHART_TYPE_ROUTER = {
    # Bar charts (horizontal)
    "barChart": _render_bar,
    "clusteredBarChart": _render_bar,
    "stackedBarChart": _render_stacked_bar,
    "hundredPercentStackedBarChart": _render_stacked_bar,

    # Column charts (vertical)
    "columnChart": _render_column,
    "clusteredColumnChart": _render_column,
    "stackedColumnChart": _render_stacked_column,
    "hundredPercentStackedColumnChart": _render_stacked_column,

    # Line and area
    "lineChart": _render_line,
    "areaChart": _render_area,
    "stackedAreaChart": _render_area,

    # Pie and donut
    "pieChart": _render_pie,
    "donutChart": _render_donut,

    # Scatter
    "scatterChart": _render_scatter,

    # Waterfall
    "waterfallChart": _render_waterfall,

    # Combo (dual axis)
    "lineStackedColumnComboChart": _render_combo,
    "lineClusteredColumnComboChart": _render_combo,

    # Funnel
    "funnelChart": _render_funnel,

    # Treemap
    "treemap": _render_treemap,

    # Gauge
    "gauge": _render_gauge,

    # Cards and KPI
    "card": _render_card,
    "cardVisual": _render_card,
    "multiRowCard": _render_card,
    "kpi": _render_kpi,

    # Tables
    "tableEx": _render_table,
    "pivotTable": _render_table,

    # Ribbon (approximate as stacked area)
    "ribbonChart": _render_ribbon,
}

# Visual types to skip (not meaningful as static chart images)
SKIP_TYPES = {
    "slicer", "advancedSlicerVisual",
    "map", "filledMap", "shapeMap", "azureMap",
    "decompositionTreeVisual", "keyDriversVisual", "qnaVisual", "aiNarratives",
    "scriptVisual", "pythonVisual", "paginator", "referenceLabel",
}


# =============================================================================
# NATIVE PPTX CHART HELPERS
# =============================================================================

def _create_presentation():
    """Create a blank 16:9 widescreen Presentation."""
    prs = _PptxFactory()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def _build_category_chart_data(df, categories, values, series=None):
    """Build CategoryChartData from DataFrame for bar/column/line/area/pie charts.

    Pivot logic priority:
    1. If `series` is provided (explicit Legend well), pivot on that column.
    2. If 2+ grouping columns and 1 measure, pivot on the second grouping (legacy).
    3. Otherwise each measure is a series.

    Returns:
        (CategoryChartData, num_series: int)
    """
    chart_data = CategoryChartData()

    # Sort category axis by calendar month order if applicable
    df = _sort_by_month(df, categories[0])

    # Well-aware: explicit Legend/Series column
    if series and len(series) == 1 and len(values) == 1 and categories:
        pivot_df = df.pivot_table(
            index=categories[0], columns=series[0],
            values=values[0], aggfunc="sum"
        ).fillna(0)
        # Restore month order after pivot (pivot_table re-sorts index)
        month_keys = {v.lower(): i for v, i in _MONTH_ORDER.items()}
        idx_lower = [str(c).lower() for c in pivot_df.index]
        if all(v in month_keys for v in idx_lower):
            pivot_df = pivot_df.iloc[sorted(range(len(pivot_df)),
                                            key=lambda i: month_keys[idx_lower[i]])]
        chart_data.categories = [str(c) for c in pivot_df.index]
        for col in pivot_df.columns:
            chart_data.add_series(str(col), pivot_df[col].tolist())
        return chart_data, len(pivot_df.columns)

    if len(categories) >= 2 and len(values) == 1:
        # Legacy: pivot second grouping column into series
        pivot_df = df.pivot_table(
            index=categories[0], columns=categories[1],
            values=values[0], aggfunc="sum"
        ).fillna(0)
        idx_lower = [str(c).lower() for c in pivot_df.index]
        if all(v in _MONTH_ORDER for v in idx_lower):
            pivot_df = pivot_df.iloc[sorted(range(len(pivot_df)),
                                            key=lambda i: _MONTH_ORDER[idx_lower[i]])]
        chart_data.categories = [str(c) for c in pivot_df.index]
        for col in pivot_df.columns:
            chart_data.add_series(str(col), pivot_df[col].tolist())
        return chart_data, len(pivot_df.columns)
    else:
        chart_data.categories = df[categories[0]].astype(str).tolist()
        for v in values:
            chart_data.add_series(v, df[v].fillna(0).tolist())
        return chart_data, len(values)


def _build_xy_chart_data(df, spec):
    """Build XyChartData for scatter charts.

    Returns:
        (XyChartData, num_series: int) or (None, 0) if insufficient data
    """
    categories, values = classify_columns(df, spec)
    chart_data = XyChartData()

    if len(values) < 2:
        return None, 0

    num_series = 0
    if categories:
        groups = df.groupby(categories[0])
        for name, group in groups:
            series = chart_data.add_series(str(name))
            for _, row in group.iterrows():
                x_val = float(row[values[0]]) if pd.notna(row[values[0]]) else 0
                y_val = float(row[values[1]]) if pd.notna(row[values[1]]) else 0
                series.add_data_point(x_val, y_val)
            num_series += 1
    else:
        series = chart_data.add_series("Data")
        for _, row in df.iterrows():
            x_val = float(row[values[0]]) if pd.notna(row[values[0]]) else 0
            y_val = float(row[values[1]]) if pd.notna(row[values[1]]) else 0
            series.add_data_point(x_val, y_val)
        num_series = 1

    return chart_data, num_series


_C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"

# Calendar order for month names (abbreviated and full)
_MONTH_ORDER = {
    m: i for i, m in enumerate([
        "jan","feb","mar","apr","may","jun",
        "jul","aug","sep","oct","nov","dec",
        "january","february","march","april","june",
        "july","august","september","october","november","december",
    ])
}

def _sort_by_month(df, col):
    """Return df sorted by calendar month order if col contains month names.
    Falls back to original order if col doesn't look like months."""
    vals = df[col].dropna().astype(str).str.lower().unique()
    if all(v in _MONTH_ORDER for v in vals):
        df = df.copy()
        df["_month_sort"] = df[col].str.lower().map(_MONTH_ORDER)
        df = df.sort_values("_month_sort").drop(columns=["_month_sort"])
    return df


def _suppress_zero_data_labels(chart, series_values_list):
    """Inject <c:dLbl><c:idx val="N"/><c:delete val="1"/></c:dLbl> for zero-value
    data points so they don't show a cluttering '0' label on the chart.

    series_values_list: list of lists, one per series, each containing numeric values
    (in the same order as the chart series). Zero/NaN values get their label deleted.
    """
    C = _C_NS
    plot_elem = chart._element.find(f".//{{{C}}}plotArea")
    if plot_elem is None:
        return
    # Collect all <c:ser> elements across all chart types in plotArea
    ser_elems = plot_elem.findall(f".//{{{C}}}ser")
    for ser_idx, ser_el in enumerate(ser_elems):
        if ser_idx >= len(series_values_list):
            break
        values = series_values_list[ser_idx]
        zero_indices = [i for i, v in enumerate(values) if v is None or v == 0]
        if not zero_indices:
            continue
        # Find or create <c:dLbls> inside this <c:ser>
        dlbls_el = ser_el.find(f"{{{C}}}dLbls")
        if dlbls_el is None:
            # Create and insert <c:dLbls> — must appear before <c:marker>/<c:invertIfNegative>
            dlbls_el = etree.SubElement(ser_el, f"{{{C}}}dLbls")
        for pt_idx in zero_indices:
            dlbl = etree.SubElement(dlbls_el, f"{{{C}}}dLbl")
            idx_el = etree.SubElement(dlbl, f"{{{C}}}idx")
            idx_el.set("val", str(pt_idx))
            del_el = etree.SubElement(dlbl, f"{{{C}}}delete")
            del_el.set("val", "1")


def _set_dlbl_pos_xml(data_labels, pos_val):
    """Inject <c:dLblPos val="..."/> directly into a DataLabels XML element.

    The python-pptx .position API inserts a <c:txPr> sibling that causes
    PowerPoint to reject pie/donut chart files. Direct XML injection avoids
    this by placing <c:dLblPos> before <c:showLegendKey> per the OOXML schema.
    """
    dl_elem = data_labels._element
    tag = f"{{{_C_NS}}}dLblPos"
    existing = dl_elem.find(tag)
    if existing is not None:
        existing.set("val", pos_val)
        return
    pos_elem = etree.SubElement(dl_elem, tag)
    pos_elem.set("val", pos_val)
    # Move it before showLegendKey to match OOXML sequence
    show_key = dl_elem.find(f"{{{_C_NS}}}showLegendKey")
    if show_key is not None:
        show_key.addprevious(pos_elem)


def _style_native_chart(chart, spec, num_series, categories=None, values=None,
                        show_title=True):
    """Apply PBI styling to a native python-pptx chart.

    Sets title, legend, series colors, axis formatting, and axis labels
    to match PBI defaults.

    Args:
        chart: python-pptx Chart object
        spec: VisualSpec
        num_series: number of data series (for legend logic)
        categories: list of category column names (for axis labels)
        values: list of measure column names (for axis labels)
        show_title: whether to show chart title (default True)
    """
    categories = categories or []
    values = values or []
    is_pie_donut = spec.visual_type in ("pieChart", "donutChart")

    # Title
    if show_title:
        chart.has_title = True
        title_para = chart.chart_title.text_frame.paragraphs[0]
        title_para.text = spec.visual_name
        title_para.font.size = Pt(18)
        title_para.font.name = PBI_FONT
        title_para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    else:
        chart.has_title = False

    # Legend: always show for pie/donut (categories), show for multi-series others
    show_legend = is_pie_donut or num_series > 1
    if show_legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.name = PBI_FONT
    else:
        chart.has_legend = False

    # Apply PBI colors to each series and enable data labels
    _label_pos_map = {
        "barChart": XL_LABEL_POSITION.OUTSIDE_END,
        "clusteredBarChart": XL_LABEL_POSITION.OUTSIDE_END,
        "stackedBarChart": XL_LABEL_POSITION.CENTER,
        "hundredPercentStackedBarChart": XL_LABEL_POSITION.CENTER,
        "columnChart": XL_LABEL_POSITION.OUTSIDE_END,
        "clusteredColumnChart": XL_LABEL_POSITION.OUTSIDE_END,
        "stackedColumnChart": XL_LABEL_POSITION.CENTER,
        "hundredPercentStackedColumnChart": XL_LABEL_POSITION.CENTER,
        "lineChart": XL_LABEL_POSITION.ABOVE,
        "areaChart": XL_LABEL_POSITION.ABOVE,
        "stackedAreaChart": XL_LABEL_POSITION.ABOVE,
        "scatterChart": XL_LABEL_POSITION.RIGHT,
        "pieChart": XL_LABEL_POSITION.BEST_FIT,
        "donutChart": XL_LABEL_POSITION.BEST_FIT,
    }
    label_pos = _label_pos_map.get(spec.visual_type)

    plot = chart.plots[0]
    for i, series in enumerate(plot.series):
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = PBI_RGB_COLORS[i % len(PBI_RGB_COLORS)]
        # Data labels — show value for bar/column/scatter; suppress for line/area
        # (line charts with many points get very cluttered with per-point labels)
        _no_label_types = ("lineChart", "areaChart", "stackedAreaChart",
                           "lineClusteredColumnComboChart", "lineStackedColumnComboChart")
        if is_pie_donut:
            series.data_labels.show_value = False
            series.data_labels.show_percentage = True
            series.data_labels.show_category_name = False
        elif spec.visual_type in _no_label_types:
            series.data_labels.show_value = False
        else:
            series.data_labels.show_value = True
        if not is_pie_donut and label_pos is not None:
            try:
                series.data_labels.position = label_pos
            except Exception:
                pass

    # Axis font styling and labels
    try:
        cat_ax = chart.category_axis
        cat_ax.tick_labels.font.size = Pt(10)
        cat_ax.tick_labels.font.name = PBI_FONT
        cat_ax.tick_labels.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        # Category axis label (X-axis for column/line, Y-axis for bar)
        if categories:
            cat_ax.has_title = True
            cat_ax.axis_title.text_frame.paragraphs[0].text = categories[0]
            cat_ax.axis_title.text_frame.paragraphs[0].font.size = Pt(12)
            cat_ax.axis_title.text_frame.paragraphs[0].font.name = PBI_FONT
            cat_ax.axis_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    except Exception:
        pass  # pie/donut charts have no category axis

    try:
        val_ax = chart.value_axis
        val_ax.tick_labels.font.size = Pt(10)
        val_ax.tick_labels.font.name = PBI_FONT
        val_ax.tick_labels.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        val_ax.has_major_gridlines = True
        val_ax.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        # Auto-detect percentage format: if measure name contains % / YoY / rate / pct
        # or if all numeric values are in (-2, 2) range (likely a ratio/percentage)
        _pct_keywords = ("% ", "%", "yoy", "rate", "pct", "ratio", "change", "growth", "variance")
        _is_pct = any(kw in v.lower() for v in values for kw in _pct_keywords)
        if not _is_pct and values:
            try:
                numeric_vals = df[values[0]].dropna()
                if len(numeric_vals) > 0 and numeric_vals.abs().max() <= 20 and numeric_vals.abs().max() <= 5:
                    _is_pct = True
            except Exception:
                pass
        if _is_pct:
            val_ax.tick_labels.number_format = "0%"
            val_ax.tick_labels.number_format_is_linked = False
        # Value axis label
        if len(values) == 1:
            val_ax.has_title = True
            val_ax.axis_title.text_frame.paragraphs[0].text = values[0]
            val_ax.axis_title.text_frame.paragraphs[0].font.size = Pt(12)
            val_ax.axis_title.text_frame.paragraphs[0].font.name = PBI_FONT
            val_ax.axis_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    except Exception:
        pass  # pie/donut charts have no value axis


def _add_native_chart(slide, df, spec, left=None, top=None, width=None, height=None,
                      show_title=True):
    """Add a native python-pptx chart to the slide.

    Routes to the correct chart type based on spec.visual_type.
    Handles bar, column, line, area, pie, donut, and scatter natively.
    Optional position params override the default centered placement.

    Returns:
        True if chart was added successfully, False otherwise
    """
    chart_type_enum = NATIVE_CHART_MAP.get(spec.visual_type)
    if chart_type_enum is None:
        return False

    c_left = left if left is not None else CHART_LEFT
    c_top = top if top is not None else CHART_TOP
    c_width = width if width is not None else CHART_WIDTH
    c_height = height if height is not None else CHART_HEIGHT

    categories, values = classify_columns(df, spec)

    # --- Scatter/XY charts use XyChartData ---
    if spec.visual_type == "scatterChart":
        chart_data, num_series = _build_xy_chart_data(df, spec)
        if chart_data is None:
            return False
        chart_frame = slide.shapes.add_chart(
            chart_type_enum, c_left, c_top, c_width, c_height,
            chart_data
        )
        chart = chart_frame.chart
        # Scatter: X-axis = values[0], Y-axis = values[1]
        x_label = [values[0]] if values else []
        y_label = [values[1]] if len(values) >= 2 else []
        _style_native_chart(chart, spec, num_series,
                            categories=x_label, values=y_label,
                            show_title=show_title)
        # Scatter markers: set size
        for series in chart.plots[0].series:
            series.marker.size = 10
        return True

    # --- Pie/donut: measures-only case (each measure is a slice) ---
    if spec.visual_type in ("pieChart", "donutChart"):
        if not categories and len(values) >= 2:
            chart_data = CategoryChartData()
            chart_data.categories = values
            row = df.iloc[0]
            chart_data.add_series(
                "Values",
                [float(row[v]) if pd.notna(row[v]) else 0 for v in values]
            )
            chart_frame = slide.shapes.add_chart(
                chart_type_enum, c_left, c_top, c_width, c_height,
                chart_data
            )
            chart = chart_frame.chart
            _style_native_chart(chart, spec, 1, show_title=show_title)
            # Color individual pie/donut slices
            plot = chart.plots[0]
            for i, point in enumerate(plot.series[0].points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = PBI_RGB_COLORS[i % len(PBI_RGB_COLORS)]
            return True

    if not categories or not values:
        return False

    # Sort bar/column charts by first measure descending (matches PBI default)
    if spec.visual_type in ("barChart", "clusteredBarChart",
                            "columnChart", "clusteredColumnChart"):
        if len(values) == 1 and len(categories) == 1:
            df = df.sort_values(values[0], ascending=False)

    # --- Standard category charts (bar, column, line, area, pie, donut) ---
    series_cols = _resolve_series(df, spec)
    chart_data, num_series = _build_category_chart_data(df, categories, values, series_cols)
    chart_frame = slide.shapes.add_chart(
        chart_type_enum, c_left, c_top, c_width, c_height,
        chart_data
    )
    chart = chart_frame.chart
    _style_native_chart(chart, spec, num_series,
                        categories=categories, values=values,
                        show_title=show_title)

    # Suppress "0" data labels on zero-fill pivot entries (clutters clustered charts)
    if num_series > 1:
        series_vals = [list(s.values) for s in chart.plots[0].series]
        _suppress_zero_data_labels(chart, series_vals)

    # Pie/donut: color individual slices instead of series
    if spec.visual_type in ("pieChart", "donutChart"):
        plot = chart.plots[0]
        for i, point in enumerate(plot.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = PBI_RGB_COLORS[i % len(PBI_RGB_COLORS)]

    return True


def _add_png_to_slide(slide, png_path, left=None, top=None, width=None, height=None):
    """Insert a plotly PNG image onto a slide. Optional position params override defaults."""
    slide.shapes.add_picture(
        str(png_path),
        left if left is not None else CHART_LEFT,
        top if top is not None else CHART_TOP,
        width if width is not None else CHART_WIDTH,
        height if height is not None else CHART_HEIGHT,
    )


# =============================================================================
# NATIVE PPTX RENDERERS — Table, Card, KPI, Ribbon, Combo
# =============================================================================

def _format_cell_value(value, col_name=""):
    """Format a cell value for display in a PowerPoint table or card.

    Detects currency and percentage patterns from column name keywords
    and applies human-readable formatting.
    """
    if pd.isna(value) or str(value).strip() in ("", "(Blank)"):
        return ""
    col_lower = col_name.lower()
    is_currency = any(kw in col_lower for kw in (
        "revenue", "sales", "profit", "cost", "price", "amount", "spend",
        "budget", "loss", "premium", "income", "margin$", "dollar",
    ))
    is_pct = any(kw in col_lower for kw in (
        "ratio", "percent", "%", "rate", "proportion", "share",
    ))
    try:
        num = float(value)
        if is_pct and abs(num) <= 10:
            # Looks like a decimal ratio (0.85) — convert to percent
            return f"{num * 100:.1f}%"
        elif is_pct:
            return f"{num:.1f}%"
        elif is_currency:
            if abs(num) >= 1_000_000_000:
                return f"${num / 1_000_000_000:,.1f}B"
            elif abs(num) >= 1_000_000:
                return f"${num / 1_000_000:,.1f}M"
            elif abs(num) >= 1_000:
                return f"${num:,.0f}"
            else:
                return f"${num:,.2f}"
        elif num == int(num) and abs(num) < 1e15:
            return f"{int(num):,}"
        else:
            return f"{num:,.2f}"
    except (ValueError, TypeError):
        return str(value)


def _add_native_table(slide, df, spec, left=None, top=None, width=None, height=None,
                      show_title=True):
    """Add a native editable PowerPoint table to the slide.

    Styled with PBI dark-blue header, alternating row colors, auto-fit column
    widths, and Segoe UI font throughout. Caps at 50 rows.

    Returns:
        True if table was added successfully.
    """
    max_rows = 50
    display_df = df.head(max_rows)
    num_rows = len(display_df)
    num_cols = len(display_df.columns)

    if num_rows == 0 or num_cols == 0:
        return False

    # Reserve space for title
    title_height = Inches(0.5) if show_title else Inches(0)
    t_left = left if left is not None else CHART_LEFT
    t_top = (top if top is not None else CHART_TOP) + title_height
    t_width = width if width is not None else CHART_WIDTH
    t_height = (height if height is not None else CHART_HEIGHT) - title_height

    # Add title text box
    if show_title:
        title_box = slide.shapes.add_textbox(
            t_left, t_top - title_height, t_width, title_height
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        title_text = spec.visual_name
        if len(df) > max_rows:
            title_text += f" (showing {max_rows} of {len(df)} rows)"
        p.text = title_text
        p.font.size = Pt(18)
        p.font.name = PBI_FONT
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # +1 for header row
    table_shape = slide.shapes.add_table(
        num_rows + 1, num_cols, t_left, t_top, t_width, t_height
    )
    table = table_shape.table

    # Distribute column widths evenly
    col_width = int(t_width / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_width

    # Header row — dark blue background, white text
    for j, col_name in enumerate(display_df.columns):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.name = PBI_FONT
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        # Dark blue fill
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = PBI_RGB_COLORS[1]  # #12239E

    # Data rows — alternating white/#F5F5F5
    for i in range(num_rows):
        row_color = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else RGBColor(0xF5, 0xF5, 0xF5)
        for j, col_name in enumerate(display_df.columns):
            cell = table.cell(i + 1, j)
            raw_val = display_df.iloc[i, j]
            cell.text = _format_cell_value(raw_val, col_name)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.font.name = PBI_FONT
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            # Row fill
            cell_fill = cell.fill
            cell_fill.solid()
            cell_fill.fore_color.rgb = row_color

    return True


def _add_native_card(slide, df, spec, left=None, top=None, width=None, height=None,
                     show_title=True):
    """Add native PowerPoint card visual — large styled text boxes.

    Single card: one big centered number with label below.
    Multi-card: side-by-side cards evenly distributed.

    Returns:
        True if card was added successfully.
    """
    _, values = classify_columns(df, spec)
    if not values or df.empty:
        return False

    row = df.iloc[0]
    c_left = left if left is not None else CHART_LEFT
    c_top = top if top is not None else CHART_TOP
    c_width = width if width is not None else CHART_WIDTH
    c_height = height if height is not None else CHART_HEIGHT

    num_cards = len(values)
    card_width = int(c_width / num_cards)
    card_padding = Inches(0.2)

    for i, v in enumerate(values):
        val = row[v]
        formatted = _format_cell_value(val, v)

        # Card bounding box
        box_left = c_left + card_width * i + card_padding
        box_width = card_width - card_padding * 2

        # Value text box — large centered number
        value_top = c_top + Inches(1.5)
        value_height = Inches(2.0)
        value_box = slide.shapes.add_textbox(box_left, value_top, box_width, value_height)
        tf = value_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = formatted
        p.alignment = PP_ALIGN.CENTER
        # Scale font size based on number of cards
        font_size = 60 if num_cards == 1 else (44 if num_cards <= 2 else 32)
        p.font.size = Pt(font_size)
        p.font.name = PBI_FONT
        p.font.color.rgb = PBI_RGB_COLORS[i % len(PBI_RGB_COLORS)]
        p.font.bold = True

        # Label text box — measure name below
        label_top = value_top + value_height
        label_height = Inches(0.6)
        label_box = slide.shapes.add_textbox(box_left, label_top, box_width, label_height)
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = v
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.name = PBI_FONT
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # Overall title
    if show_title and spec.visual_name:
        title_box = slide.shapes.add_textbox(c_left, c_top, c_width, Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = spec.visual_name
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.name = PBI_FONT
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    return True


def _add_native_kpi(slide, df, spec, left=None, top=None, width=None, height=None,
                    show_title=True):
    """Add native PowerPoint KPI visual — big number + delta indicator.

    First measure = value, second measure = reference (delta = value - reference).
    Delta shown in green (positive) or red (negative) below the main value.

    Returns:
        True if KPI was added successfully.
    """
    _, values = classify_columns(df, spec)
    if not values or df.empty:
        return False

    row = df.iloc[0]
    c_left = left if left is not None else CHART_LEFT
    c_top = top if top is not None else CHART_TOP
    c_width = width if width is not None else CHART_WIDTH
    c_height = height if height is not None else CHART_HEIGHT

    value = float(row[values[0]]) if pd.notna(row[values[0]]) else 0
    formatted_value = _format_cell_value(value, values[0])

    # Title
    if show_title and spec.visual_name:
        title_box = slide.shapes.add_textbox(c_left, c_top, c_width, Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = spec.visual_name
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.name = PBI_FONT
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Main value — large centered
    value_top = c_top + Inches(1.5)
    value_box = slide.shapes.add_textbox(c_left, value_top, c_width, Inches(2.0))
    tf = value_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = formatted_value
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54)
    p.font.name = PBI_FONT
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.font.bold = True

    # Delta indicator (if second measure exists)
    if len(values) >= 2:
        reference = float(row[values[1]]) if pd.notna(row[values[1]]) else 0
        delta = value - reference
        if reference != 0:
            delta_pct = delta / abs(reference) * 100
            delta_text = f"{'▲' if delta >= 0 else '▼'} {abs(delta_pct):.1f}%"
        else:
            delta_text = f"{'▲' if delta >= 0 else '▼'} {_format_cell_value(abs(delta), values[0])}"
        delta_color = RGBColor(0x1A, 0xAB, 0x40) if delta >= 0 else RGBColor(0xD6, 0x45, 0x50)

        delta_top = value_top + Inches(2.0)
        delta_box = slide.shapes.add_textbox(c_left, delta_top, c_width, Inches(0.8))
        tf = delta_box.text_frame
        p = tf.paragraphs[0]
        p.text = delta_text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.name = PBI_FONT
        p.font.color.rgb = delta_color
        p.font.bold = True

        # Reference label
        ref_top = delta_top + Inches(0.8)
        ref_box = slide.shapes.add_textbox(c_left, ref_top, c_width, Inches(0.5))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"vs {values[1]}: {_format_cell_value(reference, values[1])}"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.name = PBI_FONT
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    return True


def _add_native_ribbon(slide, df, spec, left=None, top=None, width=None, height=None,
                       show_title=True):
    """Add native PowerPoint ribbon chart as AREA_STACKED.

    Ribbon charts are visually equivalent to stacked area charts.
    Uses the same pivot logic as the plotly renderer.

    Returns:
        True if chart was added successfully.
    """
    categories, values = classify_columns(df, spec)
    if not categories or not values:
        return False

    c_left = left if left is not None else CHART_LEFT
    c_top = top if top is not None else CHART_TOP
    c_width = width if width is not None else CHART_WIDTH
    c_height = height if height is not None else CHART_HEIGHT

    df_sorted = df.sort_values(categories[0])
    chart_data = CategoryChartData()

    if len(categories) >= 2 and len(values) == 1:
        # Pivot second grouping into series
        pivot_df = df_sorted.pivot_table(
            index=categories[0], columns=categories[1],
            values=values[0], aggfunc="sum"
        ).fillna(0)
        chart_data.categories = [str(c) for c in pivot_df.index]
        num_series = len(pivot_df.columns)
        for col in pivot_df.columns:
            chart_data.add_series(str(col), pivot_df[col].tolist())
    else:
        chart_data.categories = df_sorted[categories[0]].astype(str).tolist()
        num_series = len(values)
        for v in values:
            chart_data.add_series(v, df_sorted[v].tolist())

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.AREA_STACKED, c_left, c_top, c_width, c_height, chart_data
    )
    chart = chart_frame.chart
    _style_native_chart(chart, spec, num_series,
                        categories=categories, values=values,
                        show_title=show_title)
    return True


def _add_native_combo(slide, df, spec, left=None, top=None, width=None, height=None,
                      show_title=True):
    """Add native PowerPoint combo chart — columns + line on secondary axis.

    Creates a column chart, then patches the XML to add a line chart overlay
    on a secondary value axis. Uses metadata y2_columns to determine which
    measures go on the line (secondary axis).

    Single-measure + two-grouping-columns pattern: pivot second grouping as
    channel series (stacked bars per X-axis category) and add a "Total" line
    series showing the sum across all channels. This handles the common PBI
    pattern of "Channel stacked by Month + Total line".

    Returns:
        True if chart was added successfully.
    """
    categories, values = classify_columns(df, spec)
    if not categories or not values:
        return False

    c_left = left if left is not None else CHART_LEFT
    c_top = top if top is not None else CHART_TOP
    c_width = width if width is not None else CHART_WIDTH
    c_height = height if height is not None else CHART_HEIGHT

    # --- Single-measure + two-grouping-column case ---
    # This is a "small multiples" layout (e.g. Channel facets × Monthly bars).
    # PowerPoint has no native small-multiples chart type, so fall back to
    # the plotly renderer which generates properly faceted subplots.
    if len(values) == 1 and len(categories) >= 2:
        return False

    # --- Multi-measure case (original logic) ---
    if len(values) < 2:
        return False

    # Resolve y2 measures (secondary axis = line)
    df_cols_lower = {c.lower().strip(): c for c in df.columns}
    y2_actual = set()
    for y2 in spec.y2_columns:
        actual = df_cols_lower.get(y2.lower().strip())
        if actual:
            y2_actual.add(actual)

    if y2_actual:
        bar_measures = [v for v in values if v not in y2_actual]
        line_measures = [v for v in values if v in y2_actual]
    elif len(values) > 1:
        bar_measures = values[:-1]
        line_measures = [values[-1]]
    else:
        bar_measures = values
        line_measures = []

    if not line_measures:
        return False  # No secondary axis — use regular column chart

    # Determine base chart type based on visual type
    is_stacked = "stacked" in spec.visual_type.lower()
    base_chart_type = XL_CHART_TYPE.COLUMN_STACKED if is_stacked else XL_CHART_TYPE.COLUMN_CLUSTERED

    # Build chart data with ALL measures (bars first, then lines)
    all_measures = bar_measures + line_measures
    chart_data = CategoryChartData()
    cat_labels = df[categories[0]].astype(str).tolist()
    chart_data.categories = cat_labels
    for m in all_measures:
        vals = df[m].fillna(0).tolist()
        chart_data.add_series(m, vals)

    # Add the chart as column initially
    chart_frame = slide.shapes.add_chart(
        base_chart_type, c_left, c_top, c_width, c_height, chart_data
    )
    chart = chart_frame.chart

    # --- XML surgery: move line measures to a lineChart on secondary axis ---
    nsmap = {
        'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }
    chart_xml = chart._chartSpace
    plot_area = chart_xml.find('.//c:plotArea', nsmap)
    bar_chart_el = plot_area.find('c:barChart', nsmap)
    if bar_chart_el is None:
        return False

    # Create a lineChart element
    line_chart_el = etree.SubElement(plot_area, '{http://schemas.openxmlformats.org/drawingml/2006/chart}lineChart')

    # Add grouping to line chart
    grouping_el = etree.SubElement(line_chart_el, '{http://schemas.openxmlformats.org/drawingml/2006/chart}grouping')
    grouping_el.set('val', 'standard')

    # Move line series from barChart to lineChart
    bar_series_list = bar_chart_el.findall('c:ser', nsmap)
    line_series_indices = list(range(len(bar_measures), len(all_measures)))

    for idx in sorted(line_series_indices, reverse=True):
        if idx < len(bar_series_list):
            ser_el = bar_series_list[idx]
            bar_chart_el.remove(ser_el)

            # Add marker to line series for visibility
            marker_el = etree.SubElement(ser_el, '{http://schemas.openxmlformats.org/drawingml/2006/chart}marker')
            symbol_el = etree.SubElement(marker_el, '{http://schemas.openxmlformats.org/drawingml/2006/chart}symbol')
            symbol_el.set('val', 'circle')
            size_el = etree.SubElement(marker_el, '{http://schemas.openxmlformats.org/drawingml/2006/chart}size')
            size_el.set('val', '5')

            line_chart_el.append(ser_el)

    # Add axis references: line chart uses secondary value axis (catAx=primary, valAx=secondary)
    # Primary axes IDs (already exist on bar chart)
    existing_cat_ax = plot_area.find('c:catAx', nsmap)
    existing_val_ax = plot_area.find('c:valAx', nsmap)
    if existing_cat_ax is None or existing_val_ax is None:
        return False

    primary_cat_id = existing_cat_ax.find('c:axId', nsmap).get('val')
    primary_val_id = existing_val_ax.find('c:axId', nsmap).get('val')

    # Secondary axis IDs
    sec_cat_id = str(int(primary_cat_id) + 100)
    sec_val_id = str(int(primary_val_id) + 100)

    # Add axId refs to line chart element
    cat_ax_ref = etree.SubElement(line_chart_el, '{http://schemas.openxmlformats.org/drawingml/2006/chart}axId')
    cat_ax_ref.set('val', sec_cat_id)
    val_ax_ref = etree.SubElement(line_chart_el, '{http://schemas.openxmlformats.org/drawingml/2006/chart}axId')
    val_ax_ref.set('val', sec_val_id)

    # Create secondary category axis (hidden, shares labels with primary)
    sec_cat_ax = etree.SubElement(plot_area, '{http://schemas.openxmlformats.org/drawingml/2006/chart}catAx')
    _ax_id = etree.SubElement(sec_cat_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}axId')
    _ax_id.set('val', sec_cat_id)
    _scaling = etree.SubElement(sec_cat_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}scaling')
    _orient = etree.SubElement(_scaling, '{http://schemas.openxmlformats.org/drawingml/2006/chart}orientation')
    _orient.set('val', 'minMax')
    _delete = etree.SubElement(sec_cat_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}delete')
    _delete.set('val', '1')  # hidden
    _ax_pos = etree.SubElement(sec_cat_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}axPos')
    _ax_pos.set('val', 'b')
    _cross_ax = etree.SubElement(sec_cat_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}crossAx')
    _cross_ax.set('val', sec_val_id)

    # Create secondary value axis (visible, on the right side)
    sec_val_ax = etree.SubElement(plot_area, '{http://schemas.openxmlformats.org/drawingml/2006/chart}valAx')
    _ax_id2 = etree.SubElement(sec_val_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}axId')
    _ax_id2.set('val', sec_val_id)
    _scaling2 = etree.SubElement(sec_val_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}scaling')
    _orient2 = etree.SubElement(_scaling2, '{http://schemas.openxmlformats.org/drawingml/2006/chart}orientation')
    _orient2.set('val', 'minMax')
    _delete2 = etree.SubElement(sec_val_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}delete')
    _delete2.set('val', '0')  # visible
    _ax_pos2 = etree.SubElement(sec_val_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}axPos')
    _ax_pos2.set('val', 'r')  # right side
    _cross_ax2 = etree.SubElement(sec_val_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}crossAx')
    _cross_ax2.set('val', sec_cat_id)
    _crosses = etree.SubElement(sec_val_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}crosses')
    _crosses.set('val', 'max')  # cross at max so it appears on right

    # Style: font on secondary axis
    _num_fmt = etree.SubElement(sec_val_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}numFmt')
    _num_fmt.set('formatCode', 'General')
    _num_fmt.set('sourceLinked', '1')
    _major_gridlines = etree.SubElement(sec_val_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}majorGridlines')

    # Apply PBI styling to the chart (title, legend, colors for bar series)
    num_series = len(all_measures)
    _style_native_chart(chart, spec, num_series,
                        categories=categories, values=bar_measures,
                        show_title=show_title)

    # Re-color line series (they were moved, so reapply from chart object)
    # The line series colors need to be set via XML since they're in a separate plot
    line_sers = line_chart_el.findall('c:ser', nsmap)
    for i, ser in enumerate(line_sers):
        color_idx = (len(bar_measures) + i) % len(PBI_RGB_COLORS)
        rgb = PBI_RGB_COLORS[color_idx]
        hex_color = f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

        # Set line color via spPr
        sp_pr = ser.find('c:spPr', nsmap)
        if sp_pr is None:
            sp_pr = etree.SubElement(ser, '{http://schemas.openxmlformats.org/drawingml/2006/chart}spPr')
        ln = etree.SubElement(sp_pr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
        ln.set('w', '25400')  # 2pt line width
        solid_fill = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        srgb = etree.SubElement(solid_fill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        srgb.set('val', hex_color)

    return True


# Set of visual types that have native PPTX renderers (beyond NATIVE_CHART_MAP)
NATIVE_EXTENDED_MAP = {
    "tableEx": _add_native_table,
    "pivotTable": _add_native_table,
    "card": _add_native_card,
    "cardVisual": _add_native_card,
    "multiRowCard": _add_native_card,
    "kpi": _add_native_kpi,
    "ribbonChart": _add_native_ribbon,
    "lineClusteredColumnComboChart": _add_native_combo,
    "lineStackedColumnComboChart": _add_native_combo,
}


# =============================================================================
# CORE API
# =============================================================================

def _build_spec(spec, visual_type, visual_name, grouping_columns,
                measure_columns, y2_columns, page_name):
    """Build a VisualSpec from either an existing spec or individual params."""
    if spec is not None:
        return spec
    if not visual_type:
        raise ValueError("Either spec or visual_type must be provided")
    return VisualSpec(
        page_name=page_name or "",
        visual_name=visual_name or "",
        visual_type=visual_type,
        grouping_columns=grouping_columns or [],
        measure_columns=measure_columns or [],
        y2_columns=y2_columns or [],
    )


def generate_chart(df, spec=None, visual_type=None, visual_name=None,
                   grouping_columns=None, measure_columns=None,
                   y2_columns=None, page_name=""):
    """Generate a plotly Figure for a PBI visual from tabular data.

    This is the programmatic API for plotly/PNG output. Always returns a
    plotly Figure (or None). For PowerPoint output, use generate_chart_pptx().

    Args:
        df: DataFrame with DAX query results
        spec: VisualSpec (optional, takes precedence over individual params)
        visual_type: PBI visual type identifier (e.g., "barChart")
        visual_name: Chart title
        grouping_columns: list of category column names
        measure_columns: list of value/measure column names
        y2_columns: list of secondary-axis measure column names
        page_name: PBI page name (informational)

    Returns:
        plotly Figure object, or None if visual type should be skipped
    """
    spec = _build_spec(spec, visual_type, visual_name, grouping_columns,
                       measure_columns, y2_columns, page_name)

    if spec.visual_type in SKIP_TYPES:
        print(f"  Skipping: {spec.visual_name} ({spec.visual_type}) "
              f"-- not meaningful as static chart")
        return None

    if df is None or df.empty:
        print(f"  Skipping: {spec.visual_name} -- no data")
        return None

    renderer = CHART_TYPE_ROUTER.get(spec.visual_type)
    if renderer is None:
        print(f"  WARNING: Unknown visual type '{spec.visual_type}' for "
              f"'{spec.visual_name}' -- rendering as table fallback")
        renderer = _render_table
    try:
        fig = renderer(df, spec)
        print(f"  Generated: {spec.visual_name} ({spec.visual_type})")
        return fig
    except Exception as e:
        print(f"  ERROR generating chart for '{spec.visual_name}': {e}")
        return None


def generate_chart_pptx(df, spec=None, visual_type=None, visual_name=None,
                        grouping_columns=None, measure_columns=None,
                        y2_columns=None, page_name=""):
    """Generate a single-slide PowerPoint for a PBI visual from tabular data.

    Uses a native python-pptx chart when the visual type is supported (bar,
    column, line, area, pie, donut, scatter). Falls back to rendering a plotly
    PNG and embedding it as a picture on the slide for all other types.

    Args:
        df: DataFrame with DAX query results
        spec: VisualSpec (optional, takes precedence over individual params)
        visual_type: PBI visual type identifier (e.g., "barChart")
        visual_name: Chart title
        grouping_columns: list of category column names
        measure_columns: list of value/measure column names
        y2_columns: list of secondary-axis measure column names
        page_name: PBI page name (informational)

    Returns:
        pptx.presentation.Presentation (one slide), or None if skipped
    """
    spec = _build_spec(spec, visual_type, visual_name, grouping_columns,
                       measure_columns, y2_columns, page_name)

    if spec.visual_type in SKIP_TYPES:
        print(f"  Skipping: {spec.visual_name} ({spec.visual_type}) "
              f"-- not meaningful as static chart")
        return None

    if df is None or df.empty:
        print(f"  Skipping: {spec.visual_name} -- no data")
        return None

    try:
        prs = _create_presentation()
        blank_layout = prs.slide_layouts[6]  # blank slide layout
        slide = prs.slides.add_slide(blank_layout)

        # Try native chart first for standard chart types (bar, column, line, etc.)
        if spec.visual_type in NATIVE_CHART_MAP:
            success = _add_native_chart(slide, df, spec)
            if success:
                print(f"  Generated (native PPTX): {spec.visual_name} ({spec.visual_type})")
                return prs
            print(f"  Native chart failed for '{spec.visual_name}', trying extended...")

        # Try extended native renderers (table, card, KPI, ribbon, combo)
        if spec.visual_type in NATIVE_EXTENDED_MAP:
            renderer_fn = NATIVE_EXTENDED_MAP[spec.visual_type]
            success = renderer_fn(slide, df, spec)
            if success:
                print(f"  Generated (native PPTX): {spec.visual_name} ({spec.visual_type})")
                return prs
            print(f"  Extended native failed for '{spec.visual_name}', using PNG fallback")

        # PNG fallback: render with plotly, insert image on slide
        renderer = CHART_TYPE_ROUTER.get(spec.visual_type)
        if renderer is None:
            print(f"  WARNING: Unknown visual type '{spec.visual_type}' for "
                  f"'{spec.visual_name}' -- rendering as table fallback")
            renderer = _render_table

        fig = renderer(df, spec)
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp_path = tmp.name
        try:
            fig.write_image(tmp_path, format="png", width=1100, height=500, scale=2)
            _add_png_to_slide(slide, tmp_path)
        finally:
            os.unlink(tmp_path)

        print(f"  Generated (PNG fallback on PPTX): {spec.visual_name} ({spec.visual_type})")
        return prs

    except Exception as e:
        print(f"  ERROR generating chart for '{spec.visual_name}': {e}")
        return None


def save_chart(fig, output_path, width=1100, height=500, scale=2):
    """Save a plotly Figure as a PNG image.

    Args:
        fig: plotly Figure object
        output_path: file path for the PNG (directory is created if needed)
        width: image width in pixels (before scaling)
        height: image height in pixels (before scaling)
        scale: resolution multiplier (2 = 144 DPI, 3 = 216 DPI)
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    fig.write_image(
        str(output_path),
        format="png",
        width=width,
        height=height,
        scale=scale,
    )
    print(f"  Saved: {output_path}")


def save_chart_pptx(prs, output_path):
    """Save a Presentation as a .pptx file.

    Args:
        prs: pptx.presentation.Presentation object
        output_path: file path for the .pptx (directory is created if needed)
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"  Saved: {output_path}")


# =============================================================================
# CLI
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="PBI AutoGov -- Chart Generator (Skill 5): "
                    "Generate PBI-style chart images from DAX query results."
    )

    # Required: CSV data file
    parser.add_argument("--csv", required=True,
                        help="Path to CSV file with DAX query tabular data")

    # Mode 1: Metadata-driven
    parser.add_argument("--metadata",
                        help="Path to metadata Excel from Skill 1 (pbi_report_metadata.xlsx)")
    parser.add_argument("--visual",
                        help="Visual name to match in metadata (Mode 1)")

    # Mode 2: Manual/screenshot-driven (mirrors Skill 4's interface)
    parser.add_argument("--visual-type",
                        help="PBI visual type identifier (e.g., barChart, pieChart)")
    parser.add_argument("--visual-name",
                        help="Chart title (Mode 2)")
    parser.add_argument("--field", action="append", dest="fields",
                        help="Field in 'name:role' format, repeatable "
                             "(role = measure, grouping, or y2)")

    # Output options
    parser.add_argument("--format", choices=["pptx", "png"], default="pptx",
                        dest="output_format",
                        help="Output format: 'pptx' (default) for single-slide PowerPoint "
                             "with native chart or PNG fallback, 'png' for legacy plotly image")
    parser.add_argument("--report-name",
                        help="Report name for subfolder organization (e.g., 'Revenue Opportunities' "
                             "creates output/charts/Revenue_Opportunities/)")
    parser.add_argument("--output", default="output/charts/",
                        help="Output directory for chart files (default: output/charts/)")
    parser.add_argument("--width", type=int, default=1100,
                        help="Image width in pixels for PNG mode (default: 1100)")
    parser.add_argument("--height", type=int, default=500,
                        help="Image height in pixels for PNG mode (default: 500)")
    parser.add_argument("--scale", type=int, default=2,
                        help="Resolution scale factor for PNG mode (default: 2 for 144 DPI)")

    args = parser.parse_args()

    # Load CSV data
    csv_path = Path(args.csv)
    if not csv_path.exists():
        print(f"ERROR: CSV file not found: {csv_path}")
        sys.exit(1)
    df = pd.read_csv(csv_path, encoding="utf-8-sig")
    print(f"Loaded CSV: {csv_path} ({len(df)} rows, {len(df.columns)} columns)")

    # Determine mode and build VisualSpec
    if args.metadata and args.visual:
        # Mode 1: Metadata-driven
        print(f"Mode 1: Reading visual metadata from {args.metadata}")
        spec = parse_visual_from_metadata(args.metadata, args.visual)
        if spec is None:
            sys.exit(1)
    elif args.visual_type and args.fields:
        # Mode 2: Manual CLI-driven (ad-hoc usage)
        print(f"Mode 2: Using CLI-provided visual metadata")
        grouping_cols = []
        measure_cols = []
        y2_cols = []
        for field_str in args.fields:
            if ":" not in field_str:
                print(f"ERROR: Field '{field_str}' must be in 'name:role' format "
                      f"(e.g., 'Revenue:measure')")
                sys.exit(1)
            name, role = field_str.rsplit(":", 1)
            name = name.strip()
            role = role.strip().lower()
            if role == "measure":
                measure_cols.append(name)
            elif role == "grouping":
                grouping_cols.append(name)
            elif role == "y2":
                # y2 fields are measures on the secondary axis
                measure_cols.append(name)
                y2_cols.append(name)
            else:
                print(f"WARNING: Unknown role '{role}' for field '{name}' "
                      f"-- expected 'measure', 'grouping', or 'y2'")

        spec = VisualSpec(
            page_name="",
            visual_name=args.visual_name or args.visual_type,
            visual_type=args.visual_type,
            grouping_columns=grouping_cols,
            measure_columns=measure_cols,
            y2_columns=y2_cols,
        )
    else:
        print("ERROR: Provide either --metadata + --visual (Mode 1) "
              "or --visual-type + --field (Mode 2)")
        parser.print_help()
        sys.exit(1)

    print(f"\nVisual: {spec.visual_name}")
    print(f"  Type: {spec.visual_type}")
    print(f"  Grouping columns: {spec.grouping_columns}")
    print(f"  Measure columns: {spec.measure_columns}")
    if spec.y2_columns:
        print(f"  Y2 columns (secondary axis): {spec.y2_columns}")

    # Sanitize visual name for filename: replace non-alphanumeric chars with underscore
    safe_name = re.sub(r'[^\w\-]', '_', spec.visual_name).strip('_')
    output_dir = Path(args.output)
    # Create report subfolder if --report-name provided
    if args.report_name:
        safe_report = re.sub(r'[^\w\-]', '_', args.report_name).strip('_')
        output_dir = output_dir / safe_report

    # Generate and save chart
    if args.output_format == "pptx":
        prs = generate_chart_pptx(df, spec=spec)
        if prs is None:
            print("No chart generated.")
            sys.exit(0)
        output_path = output_dir / f"{safe_name}.pptx"
        save_chart_pptx(prs, output_path)
    else:
        fig = generate_chart(df, spec=spec)
        if fig is None:
            print("No chart generated.")
            sys.exit(0)
        output_path = output_dir / f"{safe_name}.png"
        save_chart(fig, output_path, width=args.width, height=args.height, scale=args.scale)

    print(f"\nDone! Chart saved to: {output_path}")


if __name__ == "__main__":
    main()
